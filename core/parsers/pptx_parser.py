"""
PPTX Parser — python-pptx + OCR for image shapes.

Per slide:
  1. Extract text from all text-frame shapes (title, body, notes)
  2. For every image shape (shape_type == 13), run OCR if Tesseract available
  3. Merge text + OCR text into the slide's content field

FIX (concept node labeling):
  - Old code: heading = title or f"Slide {slide_num}"
    → produced concept labels like "slide 10", "slide 22" when no title shape
  - New code: heading = title or _infer_heading(body_parts) or f"slide_{slide_num}"
    → tries to extract meaningful heading from first line of body content
    → "slide_N" labels (with underscore) are now filtered by extraction_engine._valid()
      because they start with "slide" — so they never become concept nodes

OCR results are labelled "slide_N_img_M" and their confidence is stored
in the section dict so the extraction engine can weight them appropriately.
"""

import re
import uuid
import os
from datetime import datetime


# ── HEADING INFERENCE ──────────────────────────────────────────────────────────

def _infer_heading(body_parts: list) -> str:
    """
    When a slide has no title shape, try to infer a meaningful heading
    from the first line of body content.

    Rules:
    - Take the first non-empty line from the first body part
    - Accept if 3-60 chars, doesn't look like a date/number/code
    - Return empty string if nothing useful found
    """
    for part in body_parts:
        lines = [l.strip() for l in part.split("\n") if l.strip()]
        for line in lines[:3]:  # check first 3 lines only
            # Skip lines that look like dates e.g. "February 10, 2020"
            if re.match(r'^(January|February|March|April|May|June|July|August|'
                        r'September|October|November|December)\s+\d', line, re.I):
                continue
            # Skip lines that are just numbers or slide refs
            if re.match(r'^[\d\s/]+$', line):
                continue
            # Skip lines that look like course headers e.g. "MATH 1112 sec. 54 Spring 2020"
            if re.search(r'(sec\.|Spring|Fall|Summer|section)\s+\d{4}', line, re.I):
                continue
            # Skip very short or very long lines
            if len(line) < 4 or len(line) > 70:
                continue
            # Skip lines that are just bullet chars
            if line in ("●", "•", "–", "-", "▶"):
                continue
            # Skip lines starting with code/bullet chars
            if line[0] in ("●", "•", "–", "▶", "#", ">"):
                continue
            # This looks like a usable heading
            return line
    return ""


# ── MAIN PARSER ───────────────────────────────────────────────────────────────

def parse_pptx(filepath: str) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from .ocr_utils import ocr_pptx_image_shape, tesseract_status

    print(f"    OCR status: {tesseract_status()}")

    prs = Presentation(filepath)
    sections, warnings, all_bold = [], [], []
    total_ocr_shapes = 0

    for i, slide in enumerate(prs.slides):
        slide_num  = i + 1
        title      = ""
        body_parts = []
        slide_bold = []
        ocr_texts  = []
        img_idx    = 0

        for shape in slide.shapes:

            # ── Image shapes: run OCR ──────────────────────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_idx += 1
                ocr_result = ocr_pptx_image_shape(shape, slide_num, img_idx)
                if ocr_result["text"]:
                    ocr_texts.append(ocr_result["text"])
                    total_ocr_shapes += 1
                    warnings.append(
                        f"Slide {slide_num} img {img_idx}: "
                        f"OCR extracted {len(ocr_result['text'].split())} words "
                        f"(conf {ocr_result['confidence']:.0f}%)"
                    )
                else:
                    warnings.append(
                        f"Slide {slide_num} img {img_idx}: "
                        f"image skipped (no readable text or OCR unavailable)"
                    )
                continue

            # ── Text frame shapes ──────────────────────────────────────
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold and run.text.strip():
                        slide_bold.append(run.text.strip())

            if "title" in shape.name.lower():
                title = shape.text_frame.text.strip()
            else:
                t = shape.text_frame.text.strip()
                if t:
                    body_parts.append(t)

        # ── Slide notes ────────────────────────────────────────────────
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        # ── Assemble content ───────────────────────────────────────────
        content_parts = body_parts[:]
        if notes:
            content_parts.append(f"NOTES: {notes}")
        if ocr_texts:
            content_parts.append(f"OCR_TEXT: {' | '.join(ocr_texts)}")

        content  = " ".join(content_parts)
        all_bold.extend(slide_bold)

        # ── Heading — use title shape, infer from body, or fallback ────
        # NOTE: fallback is now "slide_N" (underscore) not "Slide N" (space).
        # extraction_engine._valid() rejects labels starting with "slide",
        # so "slide_N" nodes are never created as concept nodes in Neo4j.
        if title:
            heading = title
        else:
            inferred = _infer_heading(body_parts)
            heading  = inferred if inferred else f"slide_{slide_num}"

        section = {
            "section_id": str(uuid.uuid4())[:8],
            "heading":    heading,
            "content":    content,
            "location":   f"slide_{slide_num}",
            "bold_terms": list(set(slide_bold)),
        }
        if ocr_texts:
            section["ocr"] = True
            section["ocr_image_count"] = img_idx

        sections.append(section)

    if total_ocr_shapes:
        print(f"    OCR: {total_ocr_shapes} image shape(s) extracted via Tesseract")

    return _build(filepath, "PPTX", sections, all_bold, warnings)


def _build(filepath, fmt, sections, bold, warnings):
    raw = " ".join(s["heading"] + " " + s["content"] for s in sections)
    return {
        "doc_id":          str(uuid.uuid4()),
        "source_file":     os.path.basename(filepath),
        "format":          fmt,
        "title":           sections[0]["heading"] if sections else os.path.basename(filepath),
        "sections":        sections,
        "raw_text":        raw,
        "all_bold_terms":  list(set(t for t in bold if t and len(t) > 2)),
        "parse_warnings":  warnings,
        "parse_success":   len(sections) > 0,
        "parse_timestamp": datetime.now().isoformat(),
    }
